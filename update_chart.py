from pptx import Presentation
from pptx.chart.data import CategoryChartData


def updateChart(pptx_file):
    presentation = Presentation(pptx_file)

    second_slide = presentation.slides[1]

    for shape in second_slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            # new_data = [10, 20, 30, 40]  # Example data, replace with your actual data
            # chart.replace_data(new_data)
            # break  # Exit the loop once the first chart is found
            # ---define new chart data---
            chart_data = CategoryChartData()
            chart_data.categories = ['Category 1', 'Category 2', 'Category 3', 'Category 4']
            chart_data.add_series('Series 1', (19.2, 21.4, 16.7, 5.8))
            chart_data.add_series('Series 2', (9.2, 15.4, 13.7, 12.1))
            chart_data.add_series('Series 3', (10.2, 25.4, 6.7, 10.3))

            # ---replace chart data---
            chart.replace_data(chart_data)
            break

    presentation.save("updated_presentation.pptx")

if __name__ == "__main__":
    pptx_file_path = "C:\\Users\\admin\\OneDrive\\Desktop\\Presentation_example4.pptx"
    updateChart(pptx_file_path)
