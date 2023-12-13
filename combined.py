from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
import csv
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE



def update_presentation(presentation_path):
    # Open the presentation
    presentation = Presentation(presentation_path)
    
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Update the title of the first slide
    first_slide = presentation.slides[0]
    title = first_slide.shapes.title
    
    title.text = "New Title"
    
    subtitle = first_slide.placeholders[1]  # Assuming the subtitle is the second placeholder
    subtitle.text = "New Subtitle"
    
    pic = "download.png"
    
   # left_inch = slide_width - width_inch  # 1 inch from the right edge
    top_inch = 0   # 1 inch from the top
    width_inch = Inches(4)  # adjust width as needed
    height_inch = Inches(2) # adjust height as needed
    left_inch = slide_width - width_inch  # 1 inch from the right edge
    first_slide.shapes.add_picture(pic, left_inch, top_inch, width_inch, height_inch)

    second_slide = presentation.slides[1]
    
    
    # Function to read data from CSV file
    def read_csv(file_path):
        with open(file_path, 'r') as csv_file:
            reader = csv.DictReader(csv_file)
            data = {row['Category']: {key: float(value) for key, value in row.items() if key != 'Category'} for row in reader}
        return data

    # Replace 'your_csv_file_path.csv' with the actual path to your CSV file
    csv_file_path = 'chart.CSV'

    for shape in second_slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            
            if chart.has_title:
                chart.chart_title.text_frame.text = "Your New Chart Title"
            else:
                title = chart.chart_title.text_frame.add_paragraph()
                title.text = "Your New Chart Title"
                chart.has_title = True
                
            chart_data = CategoryChartData()

            # Read data from CSV file
            csv_data = read_csv(csv_file_path)

            # Extract categories and series from CSV data
            categories = list(csv_data.keys())
            
            # Extract series names dynamically from the first row of CSV file
            series_names = list(next(iter(csv_data.values())).keys())
            chart_data.categories = categories

            # Add series to chart data
            for series_name in series_names:
                series_values = [csv_data[category][series_name] for category in categories]
                chart_data.add_series(series_name, series_values)

            # Replace chart data with data from CSV
            chart.replace_data(chart_data)
            break
    
    third_slide = presentation.slides[2]

    # Update with your CSV file path
    csv_file_path = 'table.csv'
    with open(csv_file_path, 'r') as csv_file:
        csv_reader = csv.reader(csv_file)
        table_data = [row for row in csv_reader]
        
    
    csv_bullets_file_path = 'bullets.csv'
    with open(csv_bullets_file_path, 'r') as csv_file:
        csv_reader = csv.reader(csv_file)
        bullets_data = [row for row in csv_reader]

    # Iterate through shapes on the third slide
    for shape in third_slide.shapes:
        # Check if the shape has a table
        if shape.has_table:
            # Iterate through rows and columns of the table
            for i, row in enumerate(shape.table.rows):
                for j, cell in enumerate(row.cells):
                    # Replace the text in each cell with data from the CSV file
                    if i < len(table_data) and j < len(table_data[i]):
                        cell.text = table_data[i][j]
                    else:
                        cell.text = ''  # In case the CSV data is smaller than the table

                    # Set the font size to 10 points
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                    
        elif shape.has_text_frame and shape.text_frame.text != '':
            paragraphs = shape.text_frame.paragraphs
            # Iterate through paragraphs (assuming bullets are in separate paragraphs)
            for i, paragraph in enumerate(paragraphs):
                # Replace the text in each paragraph with data from the CSV file
                if i < len(bullets_data):
                    paragraph.text = bullets_data[i][0]  # Assuming one column in the CSV for bullets

                    # Keep the existing font size and styling for each run in the paragraph
                    for run in paragraph.runs:
                        # You may want to adjust this line based on your specific use case
                        run.text = bullets_data[i][0]
                else:
                    paragraph.text = ''  # In case the CSV data is smaller than the number of paragraphs
                    # Clear the text in case there is no corresponding data in the CSV for this paragraph

        

    # Save the updated presentation
    presentation.save("combined_updated_presentation.pptx")


if __name__ == "__main__":
    pptx_file_path = "C:\\Users\\admin\\OneDrive\\Desktop\\Presentation_example4.pptx"
    update_presentation(pptx_file_path)
