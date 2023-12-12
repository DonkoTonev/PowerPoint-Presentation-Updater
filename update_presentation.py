from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def update_presentation(presentation_path):
    # Open the presentation
    presentation = Presentation(presentation_path)

    # Update the title of the first slide
    first_slide = presentation.slides[0]
    title = first_slide.shapes.title
    
    title.text = "New Title"
    
    subtitle = first_slide.placeholders[1]  # Assuming the subtitle is the second placeholder
    subtitle.text = "New Subtitle"    

    # Save the updated presentation
    presentation.save("updated_presentation.pptx")


# Update the path to your presentation file
presentation_path = "C:\\Users\\admin\\OneDrive\\Desktop\\Presentation_example4.pptx"
update_presentation(presentation_path)
