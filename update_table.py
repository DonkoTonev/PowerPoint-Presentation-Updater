from pptx import Presentation
from pptx.util import Pt


def updateChart(pptx_file):
    presentation = Presentation(pptx_file)

    third_slide = presentation.slides[2]

    for shape in third_slide.shapes:
        
        # if shape.has_table:

        #     shape.table.cell(1,0).text= 'Pina Colada'
        #     break
        if shape.has_table:
            cell = shape.table.cell(1, 0)
            cell.text = 'Pina Colada'

            # Set the font size to a smaller value, e.g., 10 points
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

            break

    presentation.save("updated_presentation.pptx")

if __name__ == "__main__":
    pptx_file_path = "C:\\Users\\admin\\OneDrive\\Desktop\\Presentation_example4.pptx"
    updateChart(pptx_file_path)
