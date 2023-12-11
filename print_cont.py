from pptx import Presentation

def extract_text_from_first_slide(pptx_file):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    # Access the first slide
    second_slide = presentation.slides[1]

    # Extract text from shapes in the first slide
    text_content = []
    for shape in second_slide.shapes:
        # if hasattr(shape, "text"):
        #     text_content.append(shape.text)
        if shape.has_chart:
            # text_content.append(shape)
            chart = shape.chart
            # You've found the chart; now you can update its data
            break
    new_data = [10, 20, 30, 40]  # replace with your actual new data
    chart.replace_data(new_data)
    presentation.save("C:\\Users\\admin\\OneDrive\\Desktop\\PowerPoint\\updated_presentation.pptx")

    return text_content

if __name__ == "__main__":
    # Specify the path to your PowerPoint file
    pptx_file_path = "C:\\Users\\admin\\OneDrive\\Desktop\\Presentation_example4.pptx"

    # Extract text from the first slide
    extract_text_from_first_slide(pptx_file_path)

    # Print the extracted text content
    # for text in content:
    #     print(text)
